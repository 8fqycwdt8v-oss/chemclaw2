"""Text extraction from uploaded / synced documents.

One dispatcher, `extract_text(content, content_type)`, maps a MIME type to a
per-format handler. `SUPPORTED_CONTENT_TYPES` is the single source of truth for
which types the pipeline accepts — the upload route imports it as its allowlist
so the route and the extractor can never drift.

Optional parser libraries (`pypdf`, `python-docx`, `python-pptx`, `openpyxl`)
are imported inside their handlers so this module imports cleanly where a parser
isn't installed; a missing parser surfaces as `ExtractionError`, the same
graceful pattern the original PDF path used.
"""
from __future__ import annotations

import io
import logging
import os
import zipfile
from collections.abc import Callable

from api.agent.tool_helpers import _html_to_text

logger = logging.getLogger(__name__)


PDF = "application/pdf"
DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
HTML = "text/html"
TEXT = "text/plain"
MARKDOWN = "text/markdown"


class UnsupportedContentType(ValueError):
    """Raised when no handler is registered for a content type."""


class ExtractionError(ValueError):
    """Raised when a document cannot be parsed (corrupt file or missing parser)."""


# DOCX/PPTX/XLSX are ZIP archives; the 10 MB upload cap still admits a
# ~100:1-compressible XML payload that would expand to ~1 GB inside the
# parser. Check the declared uncompressed sizes from the central directory
# before parsing — Python's zipfile caps each member read at its declared
# size, so the declared total is a sound bound for these parsers.
_MAX_UNCOMPRESSED_BYTES = 200 * 1024 * 1024  # 200 MB


def _check_zip_expansion(content: bytes) -> None:
    try:
        with zipfile.ZipFile(io.BytesIO(content)) as zf:
            total = sum(info.file_size for info in zf.infolist())
    except zipfile.BadZipFile as e:
        raise ExtractionError("Corrupt or non-ZIP Office document") from e
    if total > _MAX_UNCOMPRESSED_BYTES:
        logger.warning(
            "zip_expansion_blocked",
            extra={"declared_bytes": total, "max_bytes": _MAX_UNCOMPRESSED_BYTES},
        )
        raise ExtractionError("Document expands beyond the allowed size")


def _extract_pdf(content: bytes) -> str:
    try:
        import pypdf
    except ImportError as e:
        raise ExtractionError("PDF support not available") from e
    reader = pypdf.PdfReader(io.BytesIO(content))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_docx(content: bytes) -> str:
    try:
        from docx import Document
    except ImportError as e:
        raise ExtractionError("DOCX support not available") from e
    _check_zip_expansion(content)
    doc = Document(io.BytesIO(content))
    parts = [p.text for p in doc.paragraphs if p.text]
    # Table cells live outside the paragraph stream.
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text for c in row.cells if c.text]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(content: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ExtractionError("PPTX support not available") from e
    _check_zip_expansion(content)
    prs = Presentation(io.BytesIO(content))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text:
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _extract_xlsx(content: bytes) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as e:
        raise ExtractionError("XLSX support not available") from e
    _check_zip_expansion(content)
    # read_only + data_only: stream cells and read cached values, not formulae.
    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    try:
        parts: list[str] = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                cells = [str(v) for v in row if v is not None]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    finally:
        wb.close()


def _extract_html(content: bytes) -> str:
    return _html_to_text(content.decode("utf-8", errors="replace"))


def _extract_plain(content: bytes) -> str:
    return content.decode("utf-8", errors="replace")


_HANDLERS: dict[str, Callable[[bytes], str]] = {
    PDF: _extract_pdf,
    DOCX: _extract_docx,
    PPTX: _extract_pptx,
    XLSX: _extract_xlsx,
    HTML: _extract_html,
    TEXT: _extract_plain,
    MARKDOWN: _extract_plain,
}

SUPPORTED_CONTENT_TYPES = frozenset(_HANDLERS)

# Content types whose files are ZIP containers (PK\x03\x04) — the route
# magic-byte-checks these the same way it checks %PDF, since the Content-Type
# header is attacker-controlled.
ZIP_CONTENT_TYPES = frozenset({DOCX, PPTX, XLSX})


_EXT_CONTENT_TYPES = {
    ".pdf": PDF,
    ".docx": DOCX,
    ".pptx": PPTX,
    ".xlsx": XLSX,
    ".html": HTML,
    ".htm": HTML,
    ".txt": TEXT,
    ".md": MARKDOWN,
    ".markdown": MARKDOWN,
}


def resolve_content_type(filename: str, declared: str | None = None) -> str | None:
    """Pick a supported content type for a file, or None if unsupported.

    Trusts a `declared` MIME type (e.g. Graph's `file.mimeType`) when it's one
    we handle; otherwise falls back to the filename extension. Used by the
    drive-sync worker, where files arrive with a name + an often-imprecise
    declared type.
    """
    if declared and declared in SUPPORTED_CONTENT_TYPES:
        return declared
    _, ext = os.path.splitext(filename.lower())
    return _EXT_CONTENT_TYPES.get(ext)


def extract_text(content: bytes, content_type: str) -> str:
    """Extract plain text from `content` according to its MIME `content_type`.

    Raises `UnsupportedContentType` for an unregistered type and
    `ExtractionError` for a corrupt file or a missing parser library.
    """
    handler = _HANDLERS.get(content_type)
    if handler is None:
        raise UnsupportedContentType(content_type)
    try:
        return handler(content)
    except (UnsupportedContentType, ExtractionError):
        raise
    except Exception as e:
        # Corrupt/malformed file — log server-side, surface a generic error.
        logger.warning("extract_text_failed content_type=%s: %s", content_type, e)
        raise ExtractionError(f"Failed to parse {content_type}") from e
