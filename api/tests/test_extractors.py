"""Unit tests for document text extraction (no DB).

Builds tiny docx/pptx/xlsx/html documents in memory and asserts the dispatcher
pulls their text back out, plus the unsupported-type and bad-bytes error paths.
"""
from __future__ import annotations

import io

import pytest

from api.integrations import extractors as ex


def test_plain_and_markdown() -> None:
    assert ex.extract_text(b"hello world", ex.TEXT) == "hello world"
    assert ex.extract_text(b"# Title\n\nbody", ex.MARKDOWN) == "# Title\n\nbody"


def test_plain_invalid_utf8_replaced() -> None:
    # Lone continuation byte → replacement char, never raises.
    out = ex.extract_text(b"ok\xff", ex.TEXT)
    assert out.startswith("ok")


def test_html_strips_tags() -> None:
    html = b"<html><body><p>Reaction <b>yield</b> 92%</p><script>x=1</script></body></html>"
    out = ex.extract_text(html, ex.HTML)
    assert "Reaction" in out and "yield" in out and "92%" in out
    assert "<p>" not in out and "x=1" not in out


def test_docx_paragraphs_and_tables() -> None:
    from docx import Document

    doc = Document()
    doc.add_paragraph("First paragraph.")
    table = doc.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Solvent"
    table.rows[0].cells[1].text = "DMF"
    buf = io.BytesIO()
    doc.save(buf)

    out = ex.extract_text(buf.getvalue(), ex.DOCX)
    assert "First paragraph." in out
    assert "Solvent | DMF" in out


def test_pptx_slide_text() -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Campaign overview"
    buf = io.BytesIO()
    prs.save(buf)

    out = ex.extract_text(buf.getvalue(), ex.PPTX)
    assert "Campaign overview" in out


def test_xlsx_cells() -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws["A1"] = "compound"
    ws["B1"] = "aspirin"
    ws["A2"] = "mw"
    ws["B2"] = 180.16
    buf = io.BytesIO()
    wb.save(buf)

    out = ex.extract_text(buf.getvalue(), ex.XLSX)
    assert "compound | aspirin" in out
    assert "aspirin" in out and "180.16" in out


def test_unsupported_content_type() -> None:
    with pytest.raises(ex.UnsupportedContentType):
        ex.extract_text(b"x", "application/octet-stream")


def test_corrupt_office_file_raises_extraction_error() -> None:
    # Valid ZIP magic would still fail to parse as a real docx; arbitrary bytes
    # surface as ExtractionError, not a raw library exception.
    with pytest.raises(ex.ExtractionError):
        ex.extract_text(b"not a real docx", ex.DOCX)


def test_supported_set_matches_handlers() -> None:
    assert ex.PDF in ex.SUPPORTED_CONTENT_TYPES
    assert ex.ZIP_CONTENT_TYPES <= ex.SUPPORTED_CONTENT_TYPES
